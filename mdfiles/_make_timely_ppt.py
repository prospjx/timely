from pptx import Presentation
from pptx.util import Pt

prs = Presentation()

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Timely"
slide.placeholders[1].text = "AI-Powered Time Assistant\n5-Slide Demo Deck"

# Slide 2: Problem + Solution
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Problem & Solution"
body = slide.shapes.placeholders[1].text_frame
body.clear()
for i, text in enumerate([
    "Problem: People have tasks, but static to-do lists do not adapt to energy, urgency, or changing schedules.",
    "Solution: Timely turns natural language into structured tasks, schedules them intelligently, and delivers a daily brief.",
    "Goal: Reduce overwhelm and improve consistency through adaptive planning."
]):
    p = body.paragraphs[0] if i == 0 else body.add_paragraph()
    p.text = text
    p.level = 0

# Slide 3: Key Features
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Key Features"
body = slide.shapes.placeholders[1].text_frame
body.clear()
for i, text in enumerate([
    "Natural-language task parsing (AI extracts title, priority, deadline, duration).",
    "Smart scheduling engine with conflict checks and human-friendly pacing.",
    "Daily brief generation with voice playback.",
    "Diagnostics and reflections for behavior-aware productivity insights.",
    "Robust fallback modes for offline/partial-service reliability."
]):
    p = body.paragraphs[0] if i == 0 else body.add_paragraph()
    p.text = text
    p.level = 0

# Slide 4: Tech Stack
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "How We Built It"
body = slide.shapes.placeholders[1].text_frame
body.clear()
for i, text in enumerate([
    "Frontend: Flutter + Riverpod",
    "Backend: FastAPI + MongoDB (Motor)",
    "AI: Gemini for task parsing and brief text generation",
    "Voice: ElevenLabs with local/device fallback options",
    "Notifications: Firebase Messaging + local actionable notifications"
]):
    p = body.paragraphs[0] if i == 0 else body.add_paragraph()
    p.text = text
    p.level = 0

# Slide 5: Demo + Next Steps
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Demo Flow & Next Steps"
body = slide.shapes.placeholders[1].text_frame
body.clear()
for i, text in enumerate([
    "Demo flow: Add task -> auto-schedule -> open Daily Brief -> play voice brief.",
    "Current value: Personalized planning with resilient UX for real-world conditions.",
    "Next steps: richer personalization, stronger analytics, and multi-device readiness.",
    "Ask: Feedback on usability, realism of scheduling, and reflection quality."
]):
    p = body.paragraphs[0] if i == 0 else body.add_paragraph()
    p.text = text
    p.level = 0

# Light formatting improvements
for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame is not None:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(24) if shape == slide.shapes.title else Pt(18)

output = r"c:\Users\prosp\Downloads\MobileAP\kairos\Timely_5_Slide_Demo.pptx"
prs.save(output)
print(output)
